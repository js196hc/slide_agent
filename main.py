# main.py
# Purpose: expose a web API that builds a .pptx deck and returns a download URL.

from fastapi import FastAPI, Request, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.openapi.utils import get_openapi
from pydantic import BaseModel, HttpUrl
from pptx import Presentation
from pathlib import Path
import uuid
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from typing import Optional, Dict, List

app = FastAPI(
    title="Slide Generator",
    version="1.0.0",
    description="Creates simple PowerPoint decks from JSON input."
)

# Allow cross-origin requests for development.
# In production, tighten allow_origins to your domain.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Public directory for generated files.
PUBLIC_DIR = Path("public")
PUBLIC_DIR.mkdir(exist_ok=True)

# Mount /public so files are served over HTTP.
app.mount("/public", StaticFiles(directory=PUBLIC_DIR), name="public")

def chunk(lst, n):
    for i in range(0, len(lst), n):
        yield lst[i:i+n]

# Request schema. Validates incoming JSON.
class SlideReq(BaseModel):
    title: str
    slides: Optional[Dict[str, List[str]]] = None  # new grouped format
    bullets: Optional[List[str]] = None            # legacy support

class SlideResp(BaseModel):
    file_url: HttpUrl
    message: str

@app.post("/create_slide", response_model=SlideResp)
def create_slide(req: SlideReq, request: Request):
    try:
        # Normalize input to a slides dict
        if req.slides is None and req.bullets is not None:
            slides_dict = {"General": req.bullets}
        elif req.slides is not None:
            slides_dict = req.slides
        else:
            raise HTTPException(status_code=422, detail="Provide either slides{section:[...]} or bullets[...]")

        # tuning knobs
        MAX_PER_SLIDE = 5  # bullets per slide
        BASE_FONT_PT = 24  # normal font size
        MID_FONT_PT = 20  # if many bullets
        SMALL_FONT_PT = 18  # if very many bullets

        prs = Presentation()

        title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[1]
        cover = prs.slides.add_slide(title_layout)
        try:
            cover.shapes.title.text = req.title
        except Exception:
            pass

        for subindustry, bullets in slides_dict.items():            # split long lists across multiple slides
            parts = list(chunk(bullets, MAX_PER_SLIDE))
            total = len(parts)

            for idx, part in enumerate(parts, start=1):
                layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)

                # title: add page marker if split
                slide_title = subindustry if total == 1 else f"{subindustry} ({idx}/{total})"
                slide.shapes.title.text = slide_title

                # find body placeholder
                body_ph = None
                for ph in slide.placeholders:
                    if ph.has_text_frame and ph != slide.shapes.title:
                        body_ph = ph
                        break
                if body_ph is None:
                    continue

                tf = body_ph.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink-to-fit if needed

                # pick font size based on bullets on THIS slide
                if len(part) > 6:
                    font_pt = SMALL_FONT_PT
                elif len(part) > 4:
                    font_pt = MID_FONT_PT
                else:
                    font_pt = BASE_FONT_PT

                # build bullets
                for i, b in enumerate(part):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = b
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(font_pt)

        fname = f"{uuid.uuid4().hex}.pptx"
        out_path = PUBLIC_DIR / fname
        prs.save(out_path)
        base = str(request.base_url).rstrip("/")
        return {
            "file_url": f"{base}/public/{fname}",
            "message": f"✅ Healthcare deck ready: [Download PowerPoint]({base}/public/{fname})"
        }

    except Exception as e:
        print("ERROR /create_slide:", repr(e))
        raise HTTPException(status_code=500, detail=str(e))

def custom_openapi():
    if app.openapi_schema:
        return app.openapi_schema
    schema = get_openapi(
        title="Slide Generator",
        version="1.0.0",
        description="Creates simple PowerPoint decks from JSON input.",
        routes=app.routes,
    )
    schema["servers"] = [
        {"url": "https://slide-agent-xs03.onrender.com", "description": "Render deployment"}
    ]
    app.openapi_schema = schema
    return app.openapi_schema

app.openapi = custom_openapi

@app.get("/health")
def health(): return {"ok": True}